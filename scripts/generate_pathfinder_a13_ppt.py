"""
Generate PathFinder AI A13 presentation from current slide headings.
Preserves slide 1 (title/team) and slide 19 (publication); updates all other body text.
"""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "A13 Final Review PPT Pathfinder AI -18th May.pptx"
OUT = ROOT / "PathFinder_AI_A13_Presentation.pptx"
BACKUP = ROOT / "A13 Final Review PPT Pathfinder AI -18th May.backup.pptx"

# shape_index -> body text (titles on slides stay in template)
SLIDES: dict[int, dict[int, str]] = {
    2: {
        0: (
            "Introduction\n"
            "Problem & Objective\n"
            "Proposed Solution Overview\n"
            "System Architecture / Design\n"
            "System Workflow\n"
            "Implementation Details\n"
            "Results & Performance\n"
            "Evaluation & Analysis\n"
            "Challenges & Solutions\n"
            "Conclusion\n"
            "Future Scope\n"
            "Publication Status"
        ),
    },
    3: {
        0: (
            "PathFinder AI is a full-stack career guidance platform for students and job seekers.\n\n"
            "What it does:\n"
            "• Maps resume/profile skills to real job requirements with explainable scores\n"
            "• Recommends careers (KNN) and ranks jobs (Model 2 similarity + Doc2Vec fallback)\n"
            "• Builds personalized learning roadmaps (Google Gemini) with gap analysis\n"
            "• Adapts roadmaps from learner feedback using a constrained contextual bandit\n\n"
            "Stack: React (UI) | FastAPI + SQLite (API/DB) | XGBoost, Doc2Vec, KNN | Gemini | ChromaDB RAG (chat)\n"
            "Roles: Job seeker (profile, match, roadmap) and Recruiter (post/manage jobs)"
        ),
    },
    4: {
        0: (
            "Problem statement:\n"
            "• Generic portals use keywords — weak explainability and no learning path\n"
            "• Users cannot see which JD skills they lack or how important each skill is\n"
            "• Static roadmaps do not react when a task is too hard or off-target\n"
            "• Seekers and recruiters lack one integrated workflow\n\n"
            "Objectives:\n"
            "• JWT-secured dual-role web application\n"
            "• Resume understanding + structured user profile\n"
            "• Explainable job ranking (similarity %, matched/missing skills)\n"
            "• Model 1: JD skill importance (core / important / supporting / optional)\n"
            "• Model 2: weighted user–job similarity matching (not fit-class labels)\n"
            "• Gemini roadmaps + contextual bandit adaptation from task feedback"
        ),
    },
    5: {
        0: (
            "End-to-end pipeline:\n\n"
            "1. Data & skills — 460-skill lexicon; competency extraction from JD text; "
            "resume skills via Gemini (local lexicon fallback on API limits)\n\n"
            "2. Model 1 (XGBoost) — labels each JD skill by importance; cached on Job until JD changes\n\n"
            "3. Model 2 (similarity) — scores user vs job 0–100 from weighted overlap, "
            "experience alignment, light lexical match; lists matched/missing/critical skills\n\n"
            "4. Doc2Vec baseline — semantic embedding + 70/30 skill hybrid when profile is sparse\n\n"
            "5. Gemini — phased roadmap JSON (role summary, gaps, tasks with jd_alignment)\n\n"
            "6. Contextual bandit (ε=0.15) — 7 roadmap actions, 10-D state, action masking; "
            "Gemini rewrites tasks only when the bandit selects an edit action"
        ),
    },
    6: {
        0: (
            "Design principles:\n"
            "• Modular matchers: semantic (Doc2Vec) + structured (Model 1 + Model 2)\n"
            "• Persist expensive JD analysis once per job (Model 1)\n"
            "• Bandit for roadmaps only — not MDP/deep RL; safe masked exploration\n"
            "• Optional RAG chat grounded in seeded career documents\n\n"
            "User journey:\n"
            "1. Register → build profile → upload resume\n"
            "2. Career recommendations (KNN) + browse/search jobs\n"
            "3. Job matching → ranked by similarity_score\n"
            "4. Generate & save roadmap per job (upsert — survives refresh)\n"
            "5. Per-task feedback: Complete | Too Hard | Too Easy | Skip & Regenerate\n"
            "6. Phase 2: log interaction → recommend action → adapt roadmap\n"
            "7. Recruiter posts jobs → Model 1 skills auto-analyzed on create"
        ),
    },
    7: {
        0: (
            "System context (C4-style):\n"
            "• Actors: Job seeker, Recruiter\n"
            "• PathFinder AI: React SPA + FastAPI API + SQLite\n"
            "• External: Google Gemini API; optional ChromaDB vector store\n\n"
            "Major containers:\n"
            "• Frontend — routing, auth, job board, matching UI, roadmap detail, chat widget\n"
            "• Backend — auth, jobs CRUD, profiles, match routes, Phase 2, roadmaps\n"
            "• ML layer — KNN careers, Doc2Vec jobs, Model 1/2 services, rl_model.pkl bandit"
        ),
    },
    8: {
        0: (
            "Layers & deployment:\n"
            "• Presentation: React SPA — http://localhost:3000\n"
            "• API: FastAPI + uvicorn — http://localhost:8001 (Swagger /docs)\n"
            "• Persistence: SQLite pathfinder.db; Alembic migrations\n\n"
            "ML / AI services:\n"
            "• KNN + MultiLabelBinarizer — career recommendations\n"
            "• gensim Doc2Vec — job embedding similarity (fallback matcher)\n"
            "• Model 1 — XGBoost skill importance (model1.pkl)\n"
            "• Model 2 — rule-based similarity scoring (uses Model 1 output at runtime)\n"
            "• Bandit — rl_model.pkl, 7×10 linear policy, ε-greedy 0.15\n"
            "• Gemini — roadmaps + optional task rewrites; ChromaDB RAG for chat\n\n"
            "Training data: datasets/ (799 jobs, 12K skill rows, skills.json lexicon)"
        ),
    },
    9: {
        0: (
            "Runtime workflow:\n\n"
            "Seeker: Login → Profile/Resume → POST /api/ai/recommend-careers (KNN)\n"
            "→ Search jobs → POST /api/ai/match-jobs (Model 2 rank, Doc2Vec fallback)\n"
            "→ Job detail → POST /api/jobs/{id}/generate-roadmap-for-user (Gemini)\n"
            "→ Save roadmap (upsert per user+job)\n\n"
            "Adaptation loop on roadmap task:\n"
            "→ POST /api/phase2/interactions/log (reward + credit assignment)\n"
            "→ GET /api/phase2/recommend (masked bandit arm + valid_actions)\n"
            "→ POST /api/phase2/roadmap/adapt (apply JSON edit; Gemini if needed)\n\n"
            "Recruiter: POST job with JD → analyze-jd persisted → seekers match against stored skills"
        ),
    },
    10: {
        2: (
            "Technologies:\n"
            "• Frontend: React, React Router, Axios, JWT in localStorage\n"
            "• Backend: FastAPI, SQLAlchemy, Alembic, Pydantic schemas, bcrypt/JWT\n"
            "• Database: SQLite (dev); auto column migrate for jd_analyzed_skills fields\n\n"
            "AI / ML libraries:\n"
            "• scikit-learn (KNN), gensim (Doc2Vec), XGBoost (Model 1 training)\n"
            "• sentence-transformers / MiniLM features in Model 1 pipeline\n"
            "• google-generativeai (Gemini with model fallback chain)\n"
            "• chromadb — RAG retrieval for chatbot\n\n"
            "DevOps: start_all.bat; seed_rag.py; e2e_adaptive_roadmap_flow.py"
        ),
    },
    11: {
        2: (
            "Implemented features:\n"
            "• Auth: /api/auth/register-user | register-recruiter | token login\n"
            "• Profile + resume upload; skill extraction\n"
            "• Careers: POST /api/ai/recommend-careers\n"
            "• Jobs: search/filter, compare (up to 3), detail, recruiter CRUD\n"
            "• Model 1: POST /analyze-jd, GET /jobs/{id}/analyzed-skills (cached)\n"
            "• Model 2: POST /match-user-job, POST /recommend-jobs; match-jobs uses Model 2 first\n"
            "• Roadmaps: generate, get, upsert save per job\n"
            "• Phase 2: interactions/log, recommend, roadmap/adapt\n"
            "• Chat: RAG + session history per page context"
        ),
    },
    12: {
        0: (
            "Training data:\n"
            "• 799 balanced job postings; 12,375 (job, competency) rows; 460-skill lexicon\n\n"
            "Model 1 — XGBoost JD skill importance (held-out test n=1,857):\n"
            "• Accuracy 98.60% | macro-F1 97.32%\n"
            "• F1 by class: core 0.97 | important 0.95 | supporting 0.99 | optional 0.99\n"
            "• Key features: context_score (41%), section_score (37%)\n\n"
            "Model 2 — similarity matching (production, not high/medium/low labels):\n"
            "• similarity_score 0–100: weighted matched vs missing skills from Model 1 tiers\n"
            "• Weights core×3, important×2, supporting×1, optional×0.5\n"
            "• Blend 72% skills + 20% experience + 8% lexical; penalty for missing core skills\n\n"
            "Roadmap bandit: ε=0.15, 7 actions, 10-dimensional state vector"
        ),
    },
    13: {
        0: (
            "Demonstration scenarios:\n"
            "1. New user — resume upload → skills on profile → KNN top careers\n"
            "2. Job board — search by title/location/skills; open job detail\n"
            "3. Job matching — list sorted by similarity_score with skill breakdown\n"
            "4. Analyze JD — core/important skills visible; cached after first run\n"
            "5. Generate roadmap — phases, gap_analysis, tasks; persists on reload\n"
            "6. Too Hard on a task — bandit picks DECREASE_DIFFICULTY or ADD_PREREQUISITE; "
            "Gemini rewrites task JSON\n"
            "7. Compare jobs side-by-side; recruiter flow creates job with auto skill analysis\n\n"
            "User-visible explainability: similarity %, tiered skills, gap text, adaptation reason"
        ),
    },
    14: {
        0: (
            "Objectives met:\n"
            "• Full-stack dual-role platform with secure API\n"
            "• Explainable matching (Model 1 tiers + Model 2 similarity lists)\n"
            "• Personalized Gemini roadmaps stored per user–job pair\n"
            "• Adaptive learning via constrained contextual bandit + masked actions\n"
            "• Integrated recruiter job management\n\n"
            "Analysis:\n"
            "• Model 1 gives strong offline accuracy for JD parsing — reduces repeated LLM cost\n"
            "• Model 2 similarity is interpretable and aligned with product (ranking, not classes)\n"
            "• Bandit is appropriate for single-step roadmap edits with immediate feedback\n"
            "• Doc2Vec retains value as semantic fallback for thin profiles"
        ),
    },
    15: {
        0: (
            "Comparison with alternatives:\n"
            "• vs keyword job boards: importance-weighted skills + learning roadmaps\n"
            "• vs embedding-only match: adds Model 1 structure + explicit missing-skill lists\n"
            "• vs static MOOC paths: feedback-driven roadmap edits with safety constraints\n\n"
            "Strengths:\n"
            "• Clear separation: Model 1 (JD), Model 2 (user fit), bandit (roadmap), Gemini (content)\n"
            "• Reproducible training pipeline under data_pipeline/ and models/\n"
            "• OpenAPI-documented REST API; modular services\n\n"
            "Limitations:\n"
            "• Gemini quota/billing required for best roadmap quality\n"
            "• Jobs stored locally (no LinkedIn/Naukri API yet)\n"
            "• Bandit does not re-rank jobs — only adapts study plans"
        ),
    },
    16: {
        0: (
            "Challenges:\n"
            "• Noisy resumes and inconsistent JD formatting across 799 training jobs\n"
            "• Gemini rate limits (HTTP 429) on free tier during demos\n"
            "• Orchestrating KNN, Doc2Vec, XGBoost, similarity rules, bandit, and RAG\n"
            "• Keeping valid roadmap JSON after adaptive edits\n"
            "• Schema drift when adding jd_analyzed_skills columns to existing DBs\n\n"
            "Solutions:\n"
            "• Weak-label + feature pipeline for Model 1; persist analysis per job fingerprint\n"
            "• Gemini model fallback list + local resume skill extraction\n"
            "• model2_service as single similarity API; Phase 2 three-step adapt flow\n"
            "• Action masking (e.g. no skip when jd_importance ≥ 0.7)\n"
            "• ensure_job_analysis_columns() on API startup; Alembic migrations"
        ),
    },
    17: {
        0: (
            "PathFinder AI integrates career discovery, explainable job matching, and adaptive learning:\n\n"
            "• Profile/resume → KNN careers → Model 2 ranked jobs (+ Doc2Vec fallback)\n"
            "• Model 1 structures every JD; Model 2 scores user alignment transparently\n"
            "• Gemini produces roadmaps; contextual bandit personalizes them from real feedback\n"
            "• Recruiters and seekers share one platform and data model\n\n"
            "Deliverable: working academic prototype with measurable Model 1 metrics and "
            "end-to-end demo path suitable for final review and viva."
        ),
    },
    18: {
        0: (
            "Future scope:\n"
            "• Cloud deploy (Docker); PostgreSQL for production\n"
            "• External job feed APIs (LinkedIn, Indeed adapters)\n"
            "• Stronger RAG: automated indexing, citation in chat answers\n"
            "• Mock interviews and automated resume scoring\n"
            "• Larger labeled datasets; A/B tests on bandit reward shaping\n"
            "• Multilingual JD/resume support\n"
            "• Deep RL only if a simulated long-horizon environment is added — "
            "bandit remains the right v1 for local roadmap edits"
        ),
    },
}

SKIP_SLIDES = {1, 19}  # title + publication unchanged


def apply(prs: Presentation) -> int:
    n = 0
    for slide_num, shape_map in SLIDES.items():
        slide = prs.slides[slide_num - 1]
        for shape_idx, text in shape_map.items():
            shape = slide.shapes[shape_idx]
            if not hasattr(shape, "text"):
                continue
            shape.text = text
            n += 1
    return n


def main() -> None:
    if not SRC.exists():
        raise SystemExit(f"Template not found: {SRC}")
    if not BACKUP.exists():
        shutil.copy2(SRC, BACKUP)
    prs = Presentation(str(SRC))
    count = apply(prs)
    prs.save(str(OUT))
    shutil.copy2(OUT, SRC)
    also = ROOT / "A13.pptx"
    shutil.copy2(OUT, also)
    print(f"Wrote {count} text blocks -> {OUT}")
    print(f"Also updated: {SRC}, {also}")
    print(f"Slides: {len(prs.slides)} | Unchanged: {sorted(SKIP_SLIDES)}")


if __name__ == "__main__":
    main()
