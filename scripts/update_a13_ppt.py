"""Update A13 Final Review PPT to match implemented PathFinder AI (May 2026)."""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "A13 Final Review PPT Pathfinder AI -18th May.pptx"
BACKUP = ROOT / "A13 Final Review PPT Pathfinder AI -18th May.backup.pptx"
OUT = SRC

# slide_index (1-based) -> shape_index -> new text
UPDATES: dict[int, dict[int, str]] = {
    3: {
        2: (
            "Career platforms often give generic advice without clear job-fit reasoning.\n"
            "PathFinder AI is a full-stack web platform (React + FastAPI + SQLite) for students and job seekers.\n"
            "Two roles: job seeker and recruiter.\n"
            "The system: KNN career suggestions; job ranking (Doc2Vec + importance-weighted similarity); "
            "Gemini roadmaps with gap analysis; contextual bandit adapts saved roadmaps from task feedback.\n"
            "Stack: React, FastAPI, SQLite, ChromaDB (optional RAG), XGBoost, gensim Doc2Vec, Google Gemini."
        ),
    },
    4: {
        0: (
            "Problem statement:\n"
            "• Keyword-only matching lacks explainability and structured learning paths.\n"
            "• User skills are not mapped clearly to real JD requirements.\n"
            "• Roadmaps are static and do not adapt when tasks are too hard or irrelevant.\n"
            "• Recruiters and seekers need one platform for jobs, matching, and guidance.\n\n"
            "Objectives:\n"
            "• Build JWT-authenticated user and recruiter workflows.\n"
            "• Resume upload with Gemini + lexicon-based skill extraction.\n"
            "• KNN career recommendations; Doc2Vec + Model 2 weighted job ranking.\n"
            "• Model 1: classify JD skills as core / important / supporting / optional (saved per job).\n"
            "• Gemini JSON roadmaps; Phase 2 contextual bandit (7 actions) on task feedback."
        ),
    },
    5: {
        0: (
            "Pipeline (as implemented):\n"
            "• Skills: datasets/skills.json lexicon + competency extractor; resume via Gemini (local fallback on quota).\n"
            "• Model 1 (XGBoost): JD skill importance — stored once per job until JD changes.\n"
            "• Model 2: similarity score (0–100) from weighted matched/missing skills + experience (no high/medium/low labels).\n"
            "• Doc2Vec baseline: semantic match + 70/30 skill overlap (fallback matcher).\n"
            "• Gemini: structured roadmaps (role summary, gap analysis, phased tasks).\n"
            "• Contextual bandit (ε=0.15, 10-D state, 7 roadmap actions) — adapts learning plans, not job lists."
        ),
    },
    6: {
        0: (
            "Key ideas:\n"
            "• Three job-matching layers: Doc2Vec (semantic), Model 1 (JD parsing), Model 2 (weighted user–job similarity).\n"
            "• Gemini + optional ChromaDB RAG for chat and context-aware guidance.\n"
            "• Contextual bandit (not full MDP RL) edits roadmaps after Complete / Too Hard / Too Easy / Skip & Regenerate.\n\n"
            "Workflow:\n"
            "1. Register → profile + optional resume.\n"
            "2. KNN career recommendations; search/browse jobs; AI job matching.\n"
            "3. Generate roadmap for a job → saved per job (persists on refresh).\n"
            "4. Task feedback → log → recommend → adapt (Phase 2 APIs).\n"
            "5. Recruiter posts jobs with JD, skills, salary, remote options."
        ),
    },
    8: {
        0: (
            "Layers\n"
            "• Presentation: React SPA (port 3000)\n"
            "• API: FastAPI (port 8001) — auth, jobs, profiles, roadmaps, AI, Phase 2, match routes\n"
            "• Database: SQLite (pathfinder.db)\n"
            "• ML/AI: KNN + Doc2Vec; Model 1 & Model 2; rl_model.pkl (bandit); Gemini with model fallbacks\n"
            "• Optional: ChromaDB RAG (seed: scripts/seed_rag.py)\n\n"
            "Data\n"
            "• datasets/skills.json, training CSVs, models/*.pkl, backend/ml_models/"
        ),
    },
    10: {
        2: (
            "Technologies:\n"
            "• Frontend: React, React Router, Axios\n"
            "• Backend: FastAPI, SQLAlchemy, Alembic, uvicorn (8001)\n"
            "• Database: SQLite\n\n"
            "AI/ML:\n"
            "• Careers: KNN + MultiLabelBinarizer\n"
            "• Model 1: XGBoost + MiniLM features (skill importance)\n"
            "• Model 2: weighted similarity scoring (Model 1 output + user skills + experience)\n"
            "• Jobs UI/fallback: gensim Doc2Vec + cosine hybrid\n"
            "• Roadmaps: Gemini | Adaptation: contextual bandit (7 actions, ε=0.15)\n"
            "• RAG: ChromaDB; chat with per-page sessions"
        ),
    },
    11: {
        2: (
            "Key features:\n"
            "• User/recruiter JWT auth; profile + resume upload\n"
            "• KNN careers; job search/filter; job compare; job detail with roadmaps\n"
            "• Model 1: POST /analyze-jd — importance-weighted skills (cached per job)\n"
            "• Model 2: similarity ranking — matched/missing/critical skills\n"
            "• Gemini roadmaps; Phase 2 log / recommend / adapt\n"
            "• AI chat with RAG + conversation history\n"
            "• Recruiter enhanced job CRUD"
        ),
    },
    12: {
        0: (
            "Datasets & setup:\n"
            "• 799 job postings; 12,375 (job, skill) rows; 460-skill lexicon\n\n"
            "Model 1 — XGBoost skill importance (test n=1,857):\n"
            "• Accuracy 98.60% | macro-F1 97.32% (P 97.82%, R 96.83%)\n"
            "• Per-class F1: core 0.97 | important 0.95 | supporting 0.99 | optional 0.99\n"
            "• Top features: context_score 41%, section_score 37%\n\n"
            "Model 2 — similarity matching (user vs job, not fit classes):\n"
            "• similarity_score 0-100 from weighted matched vs missing skills (Model 1 tiers)\n"
            "• Weights: core 3.0 | important 2.0 | supporting 1.0 | optional 0.5\n"
            "• Blend: 72% skills | 20% experience | 8% profile-JD lexical; core-miss penalty\n"
            "• Outputs: ranked jobs, matched/missing/critical skills — no high/medium/low labels\n\n"
            "Also: Doc2Vec semantic baseline; roadmap bandit (epsilon=0.15, 7 actions, 10-D state)"
        ),
    },
    13: {
        0: (
            "Live demo flow (React :3000 + FastAPI :8001):\n"
            "1. Register -> upload resume -> profile skills (Gemini + lexicon fallback)\n"
            "2. Career recommendations (KNN on resume embedding)\n"
            "3. Search Jobs -> filter SQLite board; open job detail\n"
            "4. Job Matching -> ranked list with similarity %, matched & missing skills\n"
            "   (Model 2 primary; Doc2Vec fallback if profile sparse)\n"
            "5. Generate roadmap -> JSON phases saved per job (refresh-safe upsert)\n"
            "6. Roadmap feedback: Complete / Too Hard / Too Easy / Skip\n"
            "   -> log -> recommend -> adapt (bandit action + Gemini task rewrite)\n"
            "7. Compare up to 3 jobs side-by-side; recruiter posts JD -> Model 1 skills cached\n\n"
            "Explainability shown to user: skill tiers, similarity %, gap analysis, adaptation reason"
        ),
    },
    14: {
        0: (
            "Objectives achieved:\n"
            "• Full-stack platform with dual roles and JWT security\n"
            "• Resume + Gemini skill extraction (lexicon fallback)\n"
            "• KNN careers; three-layer job matching; saved roadmaps per job\n"
            "• Phase 2 contextual bandit on roadmap task interactions\n"
            "• Recruiter job board and searchable listings\n\n"
            "Evaluation:\n"
            "• Model 1: 98.6% test accuracy on JD skill tiers; persisted per job until JD changes\n"
            "• Model 2: weighted similarity matching (0-100) with matched/missing skill lists\n"
            "• Bandit: sample-efficient roadmap edits (not deep RL, not job re-ranking)"
        ),
    },
    15: {
        0: (
            "Comparison:\n"
            "• vs keyword portals: skill-weighted fit + roadmaps + adaptation\n"
            "• vs cosine-only: Doc2Vec + importance-aware Model 1/2\n"
            "• vs static roadmaps: constrained contextual bandit (7 actions, masking rules)\n\n"
            "Strengths: modular FastAPI; reproducible training; Swagger APIs; dual roles\n\n"
            "Limitations:\n"
            "• Gemini needs API quota/billing\n"
            "• Bandit adapts roadmaps only, not job re-ranking\n"
            "• No external job-feed APIs yet; local SQLite jobs"
        ),
    },
    16: {
        0: (
            "Challenges:\n"
            "• Noisy resumes and heterogeneous JDs\n"
            "• Gemini rate limits (429) on free tier\n"
            "• Integrating KNN, Doc2Vec, XGBoost, Gemini, bandit, RAG\n"
            "• Valid roadmap JSON after adaptive edits\n\n"
            "Solutions:\n"
            "• skills.json lexicon + Model 1 feature pipeline; persist JD analysis per job\n"
            "• Gemini model fallbacks + local resume skill extraction\n"
            "• Modular services; Phase 2 APIs; action masking for high JD-importance tasks\n"
            "• DB migration helpers; e2e test scripts"
        ),
    },
    17: {
        0: (
            "PathFinder AI delivers integrated career guidance:\n"
            "• Profile/resume → careers (KNN) → jobs (weighted similarity + search)\n"
            "• Gemini learning roadmaps with gap analysis\n"
            "• Contextual bandit adapts roadmaps from learner feedback\n"
            "• Explainable outputs: similarity %, matched/missing skills, adaptation reasons\n"
            "• Academic demo ready; extensible to cloud, external job APIs, and production RAG"
        ),
    },
    18: {
        0: (
            "Future work:\n"
            "• External job portal APIs; cloud deployment\n"
            "• Production RAG indexing pipeline; multilingual support\n"
            "• Mock interviews and resume scoring (out of current scope)\n"
            "• Larger datasets; library version pinning for ML artifacts\n"
            "• Optional deep RL only if long-horizon simulation is added — bandit is intentional v1"
        ),
    },
}


def apply_updates(prs: Presentation) -> int:
    count = 0
    for slide_num, shape_map in UPDATES.items():
        slide = prs.slides[slide_num - 1]
        for shape_idx, new_text in shape_map.items():
            shape = slide.shapes[shape_idx]
            if not hasattr(shape, "text"):
                continue
            shape.text = new_text
            count += 1
    return count


def main() -> None:
    if not SRC.exists():
        raise SystemExit(f"Missing: {SRC}")
    if not BACKUP.exists():
        shutil.copy2(SRC, BACKUP)
        print(f"Backup: {BACKUP}")
    prs = Presentation(str(SRC))
    n = apply_updates(prs)
    prs.save(str(OUT))
    also = ROOT / "A13.pptx"
    if also.exists() or SRC.exists():
        shutil.copy2(OUT, also)
    print(f"Updated {n} text blocks in {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
